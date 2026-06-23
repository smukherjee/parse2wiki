#!/usr/bin/env python3
"""
Extract text from PDF, DOCX, PPTX, XLSX, and ZIP files.

Zero token cost - uses pdftotext, tesseract, pandoc, and Python libraries.
"""

import argparse
import os
import subprocess
import tempfile
import zipfile
from pathlib import Path
from typing import Optional

import pdfplumber

from .document_pipeline import build_document_markdown, write_document_markdown


def extract_pdf_text(pdf_path: str, ocr: bool = False) -> str:
    """
    Extract text from PDF using pdfplumber, with pdftotext and OCR fallbacks.

    Fallback chain:
    1. pdfplumber (best for text PDFs)
    2. pdftotext (more robust for malformed PDFs)
    3. OCR (for image-based PDFs)
    """
    result = []

    # Try pdfplumber first
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for i, page in enumerate(pdf.pages):
                text = page.extract_text()
                if text and text.strip():
                    result.append(f"## Page {i + 1}\n\n{text}\n")
                elif ocr:
                    ocr_text = extract_pdf_ocr(pdf_path, i + 1)
                    if ocr_text:
                        result.append(f"## Page {i + 1} (OCR)\n\n{ocr_text}\n")

        if result:
            return "\n".join(result)
    except Exception as e:
        # Fall through to pdftotext
        pass

    # Fallback 2: pdftotext command line
    try:
        pdftotext_result = subprocess.run(
            ["pdftotext", "-layout", pdf_path, "-"],
            capture_output=True,
            text=True,
            timeout=60
        )
        if pdftotext_result.stdout.strip():
            return f"## Extracted with pdftotext\n\n{pdftotext_result.stdout}"
    except Exception:
        pass

    # Fallback 3: OCR (if enabled)
    if ocr:
        ocr_text = extract_pdf_ocr(pdf_path)
        if ocr_text and "failed" not in ocr_text.lower():
            return f"## Extracted with OCR\n\n{ocr_text}"

    return "No text content found (PDF may be image-only or corrupted)."


def extract_pdf_ocr(pdf_path: str, page_num: Optional[int] = None) -> str:
    """Extract text from PDF using tesseract OCR."""
    try:
        from pdf2image import convert_from_path
        import pytesseract

        if page_num:
            images = convert_from_path(pdf_path, first_page=page_num, last_page=page_num, dpi=300)
        else:
            images = convert_from_path(pdf_path, dpi=300)

        result = []
        for i, image in enumerate(images):
            text = pytesseract.image_to_string(image)
            if text.strip():
                result.append(f"Page {page_num or i + 1} (OCR): {text}")

        return "\n".join(result)
    except ImportError:
        return "OCR not available (install pdf2image and pytesseract)"
    except Exception as e:
        return f"OCR failed: {e}"


def extract_docx_text(docx_path: str) -> str:
    """Extract text from DOCX using python-docx."""
    try:
        from docx import Document
        doc = Document(docx_path)

        result = []
        for i, para in enumerate(doc.paragraphs):
            if para.text.strip():
                result.append(para.text)

        return "\n\n".join(result) if result else "No text content found."
    except Exception as e:
        return f"DOCX extraction failed: {e}"


def extract_pptx_text(pptx_path: str) -> str:
    """Extract text from PPTX using python-pptx."""
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)

        result = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)

            if slide_text:
                result.append(f"## Slide {i + 1}\n\n" + "\n".join(slide_text))

        return "\n\n".join(result) if result else "No text content found."
    except Exception as e:
        return f"PPTX extraction failed: {e}"


def extract_xlsx_text(xlsx_path: str) -> str:
    """Extract tables from XLSX as markdown using openpyxl."""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(xlsx_path, data_only=True)

        result = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            rows = []
            for row in sheet.iter_rows(values_only=True):
                if any(cell is not None for cell in row):
                    cells = [str(c) if c is not None else "" for c in row]
                    rows.append("| " + " | ".join(cells) + " |")

            if rows:
                header = rows[0]
                separator = "| " + " | ".join(["---"] * len(rows[0].split(" | ") - 1)) + " |"
                result.append(f"## {sheet_name}\n\n{header}\n{separator}\n" + "\n".join(rows[1:]))

        return "\n\n".join(result) if result else "No data found."
    except Exception as e:
        return f"XLSX extraction failed: {e}"


def extract_zip_contents(zip_path: str, output_dir: str) -> list:
    """Extract ZIP contents and process nested files."""
    extracted_files = []

    with zipfile.ZipFile(zip_path, 'r') as zf:
        # Create temp directory for extraction
        with tempfile.TemporaryDirectory() as tmpdir:
            zf.extractall(tmpdir)

            # Recursively process extracted files
            for root, _, files in os.walk(tmpdir):
                for file in files:
                    file_path = os.path.join(root, file)
                    rel_path = os.path.relpath(file_path, tmpdir)

                    # Process based on extension
                    ext = Path(file).suffix.lower()
                    if ext in ['.pdf', '.docx', '.pptx', '.xlsx']:
                        content = extract_file(file_path, ocr=False)
                        output_file = write_document_markdown(output_dir, rel_path, content)
                        extracted_files.append({
                            'path': f"{zip_path}/{rel_path}",
                            'output': str(output_file),
                            'content': content,
                            'extension': ext
                        })

    return extracted_files


def extract_file(file_path: str, ocr: bool = False) -> str:
    """Extract text from a single file based on extension."""
    ext = Path(file_path).suffix.lower()

    extractors = {
        '.pdf': lambda p: extract_pdf_text(p, ocr),
        '.docx': extract_docx_text,
        '.pptx': extract_pptx_text,
        '.xlsx': extract_xlsx_text,
    }

    extractor = extractors.get(ext)
    if extractor:
        return extractor(file_path)

    return f"Unsupported file type: {ext}"


def process_directory(input_dir: str, output_dir: str, ocr: bool = False, dry_run: bool = False) -> list:
    """Process all supported files in a directory recursively."""
    input_path = Path(input_dir)
    output_path = Path(output_dir)

    if not dry_run:
        output_path.mkdir(parents=True, exist_ok=True)

    results = []

    for file_path in input_path.rglob('*'):
        if file_path.is_file():
            ext = file_path.suffix.lower()

            if ext == '.zip':
                if dry_run:
                    results.append({'path': str(file_path), 'status': 'would_extract_zip'})
                else:
                    zip_results = extract_zip_contents(str(file_path), str(output_path))
                    results.extend(zip_results)
            elif ext in ['.pdf', '.docx', '.pptx', '.xlsx']:
                if dry_run:
                    results.append({'path': str(file_path), 'status': 'would_extract'})
                else:
                    content = extract_file(str(file_path), ocr)
                    output_file = write_document_markdown(str(output_path), str(file_path), content)
                    results.append({
                        'path': str(file_path),
                        'output': str(output_file),
                        'status': 'extracted'
                    })

    return results


def main():
    parser = argparse.ArgumentParser(description="Extract text from documents")
    parser.add_argument("--input", "-i", required=True, help="Input directory or file")
    parser.add_argument("--output", "-o", help="Output directory (for directories)")
    parser.add_argument("--ocr", action="store_true", help="Enable OCR for PDFs")
    parser.add_argument("--dry-run", action="store_true", help="Show what would be processed")

    args = parser.parse_args()

    input_path = Path(args.input)

    if input_path.is_file():
        content = extract_file(str(input_path), args.ocr)
        if args.dry_run:
            print(f"Would extract: {input_path}")
        else:
            print(build_document_markdown(str(input_path), content))
    else:
        results = process_directory(str(input_path), args.output or "./extractions", args.ocr, args.dry_run)

        if args.dry_run:
            print(f"Files to process: {len(results)}")
            for r in results:
                print(f"  - {r['path']}")
        else:
            print(f"Extracted {len(results)} files")
            for r in results:
                print(f"  ✓ {r['path']} → {r.get('output', 'nested')}")


if __name__ == "__main__":
    main()
